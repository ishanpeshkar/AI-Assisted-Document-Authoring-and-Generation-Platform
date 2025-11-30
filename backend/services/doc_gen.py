from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import io

def create_docx(project):
    doc = Document()
    doc.add_heading(project.title, 0)
    
    doc.add_paragraph(f"Topic: {project.topic}")
    
    for section in project.sections:
        doc.add_heading(section.title, level=1)
        if section.content:
            doc.add_paragraph(section.content)
            
    # Save to memory buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_pptx(project):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project.title
    subtitle.text = project.topic
    
    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    for section in project.sections:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = section.title
        
        if section.content:
            tf = body_shape.text_frame
            # Simple heuristic: split by newlines for bullets
            lines = section.content.split('\n')
            for i, line in enumerate(lines):
                if line.strip():
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line.strip()
                    p.font.size = Pt(18)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
